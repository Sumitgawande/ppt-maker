from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from io import BytesIO
from typing import List, Optional

app = FastAPI(title="Structured AI Presentation Builder API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Color themes
COLOR_THEMES = {
    "blue": {"primary": RGBColor(37, 99, 235), "secondary": RGBColor(191, 219, 254), "accent": RGBColor(30, 41, 59)},
    "green": {"primary": RGBColor(34, 197, 94), "secondary": RGBColor(187, 247, 208), "accent": RGBColor(15, 23, 42)},
    "purple": {"primary": RGBColor(147, 51, 234), "secondary": RGBColor(233, 213, 255), "accent": RGBColor(30, 41, 59)},
    "orange": {"primary": RGBColor(249, 115, 22), "secondary": RGBColor(254, 215, 170), "accent": RGBColor(30, 41, 59)},
}

class StructuredInput(BaseModel):
    topic: str
    presentation_type: str  # pitch_deck, report, sales, training
    audience: str
    goal: str
    key_points: List[str]
    theme: str = "blue"

class SlideContent(BaseModel):
    title: str
    bullets: List[str]
    notes: Optional[str] = ""

class PresentationData(BaseModel):
    slides: List[SlideContent]
    theme: str

class RegenerateRequest(BaseModel):
    slide_index: int
    regenerate_type: str  # "title", "bullets", "slide"
    current_title: Optional[str] = None
    current_bullets: Optional[List[str]] = None

def generate_slide_content(topic: str, presentation_type: str, audience: str, goal: str, key_points: List[str], slide_index: int) -> SlideContent:
    """Generate structured slide content based on input parameters."""
    slide_templates = {
        0: {
            "title": f"Introduction to {topic}",
            "bullets": [
                f"Overview of {topic} and its importance",
                f"Target audience: {audience}",
                f"Presentation goal: {goal}",
                f"Key focus areas: {', '.join(key_points[:3])}"
            ]
        },
        1: {
            "title": f"Problem Statement",
            "bullets": [
                f"Current challenges in {topic}",
                f"Impact on {audience}",
                f"Market gaps and opportunities",
                f"Business case for addressing these issues"
            ]
        },
        2: {
            "title": f"Solution Overview",
            "bullets": [
                f"Our approach to {topic}",
                f"Key benefits for {audience}",
                f"How we achieve {goal}",
                f"Unique value proposition"
            ]
        },
        3: {
            "title": f"Key Features",
            "bullets": key_points[:4] if len(key_points) >= 4 else key_points + ["Additional feature 1", "Additional feature 2"]
        },
        4: {
            "title": f"Implementation Strategy",
            "bullets": [
                f"Step-by-step approach to {topic}",
                f"Timeline and milestones",
                f"Resource requirements",
                f"Success metrics"
            ]
        },
        5: {
            "title": f"Results & Benefits",
            "bullets": [
                f"Expected outcomes for {audience}",
                f"ROI and measurable impact",
                f"Competitive advantages",
                f"Long-term value"
            ]
        },
        6: {
            "title": f"Next Steps",
            "bullets": [
                f"Action items and timeline",
                f"Contact information",
                f"Questions and discussion",
                f"Thank you for your attention"
            ]
        }
    }

    template = slide_templates.get(slide_index, slide_templates[0])
    return SlideContent(**template)

def regenerate_slide_section(slide_content: SlideContent, regenerate_type: str, topic: str, presentation_type: str) -> SlideContent:
    """Regenerate specific parts of a slide."""
    if regenerate_type == "title":
        # Generate new title based on bullets
        new_title = f"Enhanced: {slide_content.bullets[0][:50]}..." if slide_content.bullets else f"Slide about {topic}"
        return SlideContent(title=new_title, bullets=slide_content.bullets, notes=slide_content.notes)
    elif regenerate_type == "bullets":
        # Generate new bullets keeping the title
        new_bullets = [
            f"Improved point 1: {slide_content.bullets[0] if slide_content.bullets else 'Key insight'}",
            f"Enhanced point 2: {slide_content.bullets[1] if len(slide_content.bullets) > 1 else 'Additional detail'}",
            f"Advanced point 3: {slide_content.bullets[2] if len(slide_content.bullets) > 2 else 'Further explanation'}",
            f"Strategic point 4: {slide_content.bullets[3] if len(slide_content.bullets) > 3 else 'Conclusion'}"
        ]
        return SlideContent(title=slide_content.title, bullets=new_bullets, notes=slide_content.notes)
    elif regenerate_type == "slide":
        # Regenerate entire slide
        return generate_slide_content(topic, presentation_type, "", "", [], 0)  # Simplified
    return slide_content

@app.post("/generate-presentation")
def generate_presentation(input_data: StructuredInput):
    if not input_data.topic.strip():
        raise HTTPException(status_code=400, detail="Topic is required.")

    slides = []
    for i in range(7):  # Generate 7 slides
        slide_content = generate_slide_content(
            input_data.topic,
            input_data.presentation_type,
            input_data.audience,
            input_data.goal,
            input_data.key_points,
            i
        )
        slides.append(slide_content)

    return {"slides": slides, "theme": input_data.theme}

@app.post("/regenerate-section")
def regenerate_section(request: RegenerateRequest):
    # Create a SlideContent object from the current slide data
    current_slide = SlideContent(
        title=request.current_title or "Sample Title",
        bullets=request.current_bullets or ["Bullet 1", "Bullet 2"],
        notes=""
    )

    # Regenerate the specified section
    new_slide = regenerate_slide_section(current_slide, request.regenerate_type, "Presentation", "business")
    return new_slide

@app.post("/export-ppt")
def export_ppt(presentation_data: PresentationData):
    theme = COLOR_THEMES.get(presentation_data.theme, COLOR_THEMES["blue"])
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)

    # Title slide
    title_slide_layout = presentation.slide_layouts[6]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["primary"]

    title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = presentation_data.slides[0].title if presentation_data.slides else "AI Generated Presentation"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Content slides
    for slide_data in presentation_data.slides:
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        # Header bar
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = theme["primary"]
        header_shape.line.color.rgb = theme["primary"]

        # Slide title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for bullet in slide_data.bullets:
            para = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
            para.text = bullet
            para.level = 0
            para.font.size = Pt(20)
            para.font.color.rgb = theme["accent"]
            para.space_before = Pt(12)
            para.space_after = Pt(12)

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    filename = re.sub(r"[^A-Za-z0-9_-]", "_", presentation_data.slides[0].title if presentation_data.slides else "presentation")
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=\"{filename}.pptx\""},
    )
